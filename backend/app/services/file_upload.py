"""文件上传解析服务 — 校验、存储、提取文本、LLM 解析、写入 documents / communications。"""

import hashlib
import logging
import os
import uuid
from datetime import datetime

from fastapi import UploadFile
from sqlalchemy.ext.asyncio import AsyncSession

from app.config import settings
from app.models.document import Document
from app.models.communication import Communication

logger = logging.getLogger(__name__)


class FileUploadError(Exception):
    """文件上传异常。"""


class FileUploadService:
    """文件上传处理服务。"""

    def __init__(self) -> None:
        self.allowed_types = set(settings.allowed_file_types.split(","))
        self.max_size = settings.max_upload_size_mb * 1024 * 1024

    def _validate(self, filename: str, size: int) -> str:
        """校验文件类型和大小，返回扩展名。"""
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext not in self.allowed_types:
            raise FileUploadError(f"不支持的文件类型: .{ext}，允许: {', '.join(self.allowed_types)}")
        if size > self.max_size:
            raise FileUploadError(f"文件大小超过限制: {settings.max_upload_size_mb}MB")
        return ext

    async def process_upload(
        self,
        file: UploadFile,
        owner_id: str,
        db: AsyncSession,
        asset_owner_name: str | None = None,
    ) -> Document:
        """处理文件上传全流程。"""
        content_bytes = await file.read()
        file_size = len(content_bytes)
        ext = self._validate(file.filename or "unknown", file_size)

        # 1. 存储文件
        user_dir = os.path.join(settings.upload_dir, owner_id)
        os.makedirs(user_dir, exist_ok=True)
        file_id = str(uuid.uuid4())
        file_path = os.path.join(user_dir, f"{file_id}.{ext}")

        with open(file_path, "wb") as f:
            f.write(content_bytes)

        logger.info("文件已保存: %s (%d bytes)", file_path, file_size)

        # 2. 提取文本 & LLM 解析
        is_image = ext in ("png", "jpg", "jpeg", "gif", "webp", "bmp")
        # 保存原始文件名（含扩展名）
        original_filename = file.filename or "unknown"
        parsed = {"title": file.filename, "summary": None, "author": None, "tags": [], "category": None}

        if is_image and settings.llm_api_key and not settings.llm_api_key.startswith("sk-xxx"):
            # 图片文件：用视觉模型识别内容
            try:
                from app.services.llm import llm_client
                parsed = await llm_client.parse_image_file(content_bytes, ext)
                logger.info("视觉模型识别完成: %s", parsed.get("title"))
            except Exception as e:
                logger.warning("视觉模型图片解析失败: %s", e)
            text_content = parsed.get("content_text") or f"[图片文件: {ext}] {file.filename}"
        else:
            # 非图片文件：提取文本后用 LLM 解析
            text_content = self._extract_text(content_bytes, ext)
            if not text_content.strip():
                text_content = f"[{ext} 文件] {file.filename}"
            if settings.llm_api_key and not settings.llm_api_key.startswith("sk-xxx"):
                try:
                    from app.services.llm import llm_client
                    parsed = await llm_client.parse_uploaded_file(text_content, ext)
                except Exception as e:
                    logger.warning("LLM 文件解析失败: %s", e)

        # 4. 生成 Embedding
        embedding = None
        if settings.embedding_api_key and not settings.embedding_api_key.startswith("sk-xxx"):
            try:
                from app.services.llm import llm_client
                embed_text = f"{parsed.get('title', '')} {text_content}".strip()
                embedding = await llm_client.generate_embedding(embed_text[:2000])
            except Exception as e:
                logger.warning("Embedding 生成失败: %s", e)

        # 5. 写入 documents 表
        keywords_list = []
        if isinstance(parsed.get("tags"), list):
            keywords_list = parsed["tags"]

        # doc_category 必须在 CHECK 约束允许的范围内
        valid_categories = {"report", "proposal", "policy", "technical"}
        raw_category = parsed.get("category")
        doc_category = raw_category if raw_category in valid_categories else None

        now = datetime.utcnow()
        doc = Document(
            owner_id=owner_id,
            source_type="local",
            original_filename=original_filename,
            title=parsed.get("title") or file.filename,
            content_text=text_content,
            summary=parsed.get("summary"),
            author=parsed.get("author"),
            keywords=keywords_list,
            doc_category=doc_category,
            file_type=ext,
            file_size=file_size,
            file_path=file_path,
            asset_owner_name=asset_owner_name,
            feishu_created_at=now,
            feishu_updated_at=now,
            synced_at=now,
        )

        if embedding:
            doc.content_vector = embedding

        db.add(doc)
        await db.commit()
        await db.refresh(doc)

        logger.info("文件已入库: doc_id=%d, title=%s", doc.id, doc.title)
        return doc

    @staticmethod
    def _extract_text(content: bytes, ext: str) -> str:
        """从文件内容提取纯文本。"""
        if ext == "txt":
            for encoding in ("utf-8", "gbk", "latin-1"):
                try:
                    return content.decode(encoding)
                except (UnicodeDecodeError, ValueError):
                    continue
            return content.decode("utf-8", errors="replace")

        if ext == "csv":
            for encoding in ("utf-8", "gbk", "latin-1"):
                try:
                    return content.decode(encoding)
                except (UnicodeDecodeError, ValueError):
                    continue
            return content.decode("utf-8", errors="replace")

        if ext == "pdf":
            try:
                from pypdf import PdfReader
                import io
                reader = PdfReader(io.BytesIO(content))
                pages = [page.extract_text() or "" for page in reader.pages]
                return "\n".join(pages)
            except Exception as e:
                logger.warning("PDF 文本提取失败: %s", e)
                return ""

        if ext == "docx":
            try:
                from docx import Document as DocxDocument
                import io
                doc = DocxDocument(io.BytesIO(content))
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception as e:
                logger.warning("DOCX 文本提取失败: %s", e)
                return ""

        if ext in ("pptx", "ppt"):
            try:
                from pptx import Presentation
                import io
                prs = Presentation(io.BytesIO(content))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                return "\n".join(texts)
            except Exception as e:
                logger.warning("PPT 文本提取失败: %s", e)
                return ""

        if ext in ("xlsx", "xls"):
            try:
                import io
                import csv
                # 简单读取为文本
                return content.decode("utf-8", errors="replace")
            except Exception:
                return ""

        # 图片等非文本文件
        if ext in ("png", "jpg", "jpeg"):
            return f"[图片文件: {ext}]"

        return content.decode("utf-8", errors="replace")

    # ── 音频文件 → 沟通资产 ──────────────────────────────────

    AUDIO_EXTENSIONS = {"mp3", "wav", "m4a", "aac", "ogg", "flac"}

    async def _generate_title_fallback(self, content: str) -> str | None:
        """当标题缺失时，用 LLM 从内容生成标题。"""
        try:
            from app.services.llm import llm_client
            return await llm_client.generate_title_from_content(content)
        except Exception as e:
            logger.warning("LLM 生成标题失败: %s", e)
            return None

    async def _transcribe_audio(self, file_path: str) -> str:
        """调用 ASR API 将音频文件转为文字。"""
        if not settings.asr_api_key or settings.asr_api_key.startswith("sk-xxx"):
            raise FileUploadError(
                "语音转文字功能需要配置 ASR API Key。"
                "请在 .env 文件中设置 ASR_API_KEY（支持 OpenAI Whisper 或兼容接口）。"
            )

        from app.services.llm import create_openai_client

        asr_client = create_openai_client(
            api_key=settings.asr_api_key,
            base_url=settings.asr_base_url,
            timeout=300.0,
        )

        with open(file_path, "rb") as audio_file:
            response = await asr_client.audio.transcriptions.create(
                model=settings.asr_model,
                file=audio_file,
                language="zh",
            )

        transcript = response.text.strip()
        logger.info("ASR 转写完成，文字长度: %d", len(transcript))
        return transcript

    async def process_communication_upload(
        self,
        file: UploadFile,
        owner_id: str,
        db: AsyncSession,
        asset_owner_name: str | None = None,
        user_metadata: dict | None = None,
    ) -> Communication:
        """处理音频文件上传：存储 → ASR 转文字 → LLM 提取 → 写入 communications 表。"""
        content_bytes = await file.read()
        file_size = len(content_bytes)
        ext = self._validate(file.filename or "unknown", file_size)

        if ext not in self.AUDIO_EXTENSIONS:
            raise FileUploadError(f"沟通资产仅支持音频文件，不支持 .{ext}")

        meta = user_metadata or {}

        # 1. 存储文件
        user_dir = os.path.join(settings.upload_dir, owner_id)
        os.makedirs(user_dir, exist_ok=True)
        file_id = str(uuid.uuid4())
        file_path = os.path.join(user_dir, f"{file_id}.{ext}")

        with open(file_path, "wb") as f:
            f.write(content_bytes)

        logger.info("音频文件已保存: %s (%d bytes)", file_path, file_size)

        # 2. ASR 转文字
        transcript = await self._transcribe_audio(file_path)

        # 3. 构建用于 LLM 分析的上下文（用户补充 + ASR 转写）
        context_parts = []
        if meta.get("title"):
            context_parts.append(f"标题：{meta['title']}")
        if meta.get("comm_type"):
            type_labels = {"meeting": "会议录音", "phone": "电话录音", "interview": "面谈记录", "other": "其他"}
            context_parts.append(f"类型：{type_labels.get(meta['comm_type'], meta['comm_type'])}")
        if meta.get("participants"):
            context_parts.append(f"参与人：{', '.join(meta['participants'])}")
        if meta.get("context"):
            context_parts.append(f"背景信息：{meta['context']}")

        enriched_text = ""
        if context_parts:
            enriched_text = "【用户补充信息】\n" + "\n".join(context_parts) + "\n\n"
        enriched_text += "【录音转写内容】\n" + transcript

        # 4. LLM 提取结构化字段
        parsed = {
            "title": None, "comm_type": "recording", "initiator": None,
            "participants": [], "summary": None, "conclusions": None,
            "action_items": [], "keywords": [], "sentiment": "neutral",
            "duration_minutes": None, "comm_time": None,
        }
        if settings.llm_api_key and not settings.llm_api_key.startswith("sk-xxx"):
            try:
                from app.services.llm import llm_client
                parsed = await llm_client.parse_communication_doc(enriched_text)
            except Exception as e:
                logger.warning("LLM 沟通资产解析失败: %s", e)

        # 5. 生成 Embedding
        embedding = None
        if settings.embedding_api_key and not settings.embedding_api_key.startswith("sk-xxx"):
            try:
                from app.services.llm import llm_client
                embed_text = f"{parsed.get('title', '')} {transcript}".strip()
                embedding = await llm_client.generate_embedding(embed_text[:2000])
            except Exception as e:
                logger.warning("Embedding 生成失败: %s", e)

        # 6. 合并用户元数据与 LLM 提取结果（用户输入优先）
        comm_type_map = {"meeting": "meeting", "phone": "recording", "interview": "recording", "other": "recording"}
        final_comm_type = comm_type_map.get(meta.get("comm_type", ""), "") or parsed.get("comm_type") or "recording"

        participants = meta.get("participants") or parsed.get("participants") or []
        # 确保 participants 为 [{name: ...}] 格式
        if participants and isinstance(participants[0], str):
            participants = [{"name": p} for p in participants]

        comm_time = None
        if meta.get("comm_time"):
            try:
                comm_time = datetime.fromisoformat(meta["comm_time"])
            except (ValueError, TypeError):
                pass
        if not comm_time and parsed.get("comm_time"):
            try:
                comm_time = datetime.fromisoformat(parsed["comm_time"])
            except (ValueError, TypeError):
                pass

        content_hash = hashlib.md5(content_bytes).hexdigest()

        # conclusions 字段是 Text 类型，LLM 可能返回列表，需转为字符串
        raw_conclusions = parsed.get("conclusions")
        if isinstance(raw_conclusions, list):
            conclusions = "\n".join(str(c) for c in raw_conclusions)
        else:
            conclusions = raw_conclusions

        # 7. 写入 communications 表
        comm = Communication(
            owner_id=owner_id,
            comm_type=final_comm_type,
            source_platform="local",
            source_app_token=f"local_{file_id}",
            feishu_record_id=f"local_{file_id}",
            title=meta.get("title") or parsed.get("title") or await self._generate_title_fallback(transcript) or file.filename,
            comm_time=comm_time,
            initiator=parsed.get("initiator"),
            participants=participants,
            duration_minutes=parsed.get("duration_minutes"),
            conclusions=conclusions,
            action_items=parsed.get("action_items") or [],
            transcript=transcript,
            content_text=transcript,
            summary=parsed.get("summary"),
            asset_owner_name=asset_owner_name,
            keywords=parsed.get("keywords") or [],
            sentiment=parsed.get("sentiment"),
            content_hash=content_hash,
            extra_fields={
                "file_path": file_path,
                "file_size": file_size,
                "file_type": ext,
                "user_context": meta.get("context"),
            },
            parse_status="done",
            processed_at=datetime.utcnow(),
            synced_at=datetime.utcnow(),
        )

        if embedding:
            comm.content_vector = embedding

        db.add(comm)
        await db.commit()
        await db.refresh(comm)

        logger.info("沟通资产已入库: comm_id=%d, title=%s", comm.id, comm.title)
        return comm


# 模块级单例
file_upload_service = FileUploadService()
